"""
backend/genai/rag.py

RAG (Retrieval-Augmented Generation) pipeline for Decision Intelligence.
Supports: PDF, DOCX, PPTX, TXT, CSV document ingestion.

Pipeline:
  1. Upload & extract text from document
  2. Chunk text (500 chars, 50 char overlap)
  3. Generate embeddings (sentence-transformers)
  4. Store in ChromaDB (persistent)
  5. Query: embed question → retrieve top-k chunks → generate answer with citations

Routes (Blueprint prefix: /api/rag):
  POST   /api/rag/upload              — Upload & ingest document
  POST   /api/rag/query               — Query knowledge base
  GET    /api/rag/documents           — List ingested documents
  DELETE /api/rag/documents/<doc_id>  — Remove document
  DELETE /api/rag/documents           — Clear all documents
"""

import os
import io
import uuid
import json
import logging
import traceback
from pathlib import Path
from datetime import datetime

from flask import Blueprint, request, jsonify

logger = logging.getLogger(__name__)

rag_bp = Blueprint("rag", __name__)

# ── Configuration ────────────────────────────────────────────────────
CHROMA_DIR      = os.environ.get("CHROMA_PERSIST_DIR", "./chroma_db")
COLLECTION_NAME = "decision_intelligence_docs"
CHUNK_SIZE      = 500
CHUNK_OVERLAP   = 50
TOP_K           = 5     # Number of chunks to retrieve per query
MAX_FILE_MB     = int(os.environ.get("MAX_UPLOAD_SIZE_MB", "50"))
ALLOWED_EXT     = {".pdf", ".docx", ".pptx", ".txt", ".csv"}

# ── Lazy-loaded singletons ────────────────────────────────────────────
_chroma_client     = None
_chroma_collection = None
_embedder          = None


def _get_embedder():
    """Load sentence-transformers model (lazy, cached)."""
    global _embedder
    if _embedder is None:
        try:
            from sentence_transformers import SentenceTransformer
            _embedder = SentenceTransformer("all-MiniLM-L6-v2")
            logger.info("Sentence-transformers loaded: all-MiniLM-L6-v2")
        except ImportError:
            raise ImportError(
                "sentence-transformers not installed. Run: pip install sentence-transformers"
            )
    return _embedder


def _get_collection():
    """Get (or create) ChromaDB collection (lazy, cached)."""
    global _chroma_client, _chroma_collection
    if _chroma_collection is None:
        try:
            import chromadb
            _chroma_client = chromadb.PersistentClient(path=CHROMA_DIR)
            _chroma_collection = _chroma_client.get_or_create_collection(
                name=COLLECTION_NAME,
                metadata={"hnsw:space": "cosine"},
            )
            logger.info(f"ChromaDB collection ready: {COLLECTION_NAME} at {CHROMA_DIR}")
        except ImportError:
            raise ImportError(
                "chromadb not installed. Run: pip install chromadb"
            )
    return _chroma_collection


# ── Text Extraction ───────────────────────────────────────────────────

def _extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Extract plain text from uploaded file based on extension.
    Raises ValueError for unsupported types.
    """
    ext = Path(filename).suffix.lower()

    if ext == ".txt":
        return file_bytes.decode("utf-8", errors="replace")

    elif ext == ".csv":
        import csv
        text = file_bytes.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        rows = [", ".join(row) for row in reader]
        return "\n".join(rows)

    elif ext == ".pdf":
        try:
            import fitz  # PyMuPDF
            doc  = fitz.open(stream=file_bytes, filetype="pdf")
            pages = []
            for i, page in enumerate(doc):
                txt = page.get_text()
                if txt.strip():
                    pages.append(f"[Page {i+1}]\n{txt.strip()}")
            doc.close()
            return "\n\n".join(pages)
        except ImportError:
            raise ImportError("PyMuPDF not installed. Run: pip install PyMuPDF")

    elif ext == ".docx":
        try:
            from docx import Document
            doc   = Document(io.BytesIO(file_bytes))
            paras = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paras)
        except ImportError:
            raise ImportError("python-docx not installed. Run: pip install python-docx")

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs   = Presentation(io.BytesIO(file_bytes))
            texts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    texts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
            return "\n\n".join(texts)
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ── Chunking ──────────────────────────────────────────────────────────

def _chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """
    Split text into overlapping chunks for embedding.
    Prefers splitting at sentence boundaries when possible.
    """
    if not text.strip():
        return []

    chunks = []
    start  = 0
    while start < len(text):
        end   = min(start + chunk_size, len(text))
        chunk = text[start:end]
        # Try to end at a sentence boundary
        if end < len(text):
            last_period = chunk.rfind(". ")
            if last_period > chunk_size * 0.6:
                end   = start + last_period + 1
                chunk = text[start:end]
        chunks.append(chunk.strip())
        start = end - overlap
    return [c for c in chunks if c]


# ── Routes ────────────────────────────────────────────────────────────

@rag_bp.route("/upload", methods=["POST"])
def upload_document():
    """
    Upload and ingest a document into the RAG knowledge base.
    Multipart form data:
      - file: the document (PDF/DOCX/PPTX/TXT/CSV)
      - title: optional display name
    """
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    f        = request.files["file"]
    filename = f.filename or "unknown"
    ext      = Path(filename).suffix.lower()
    title    = request.form.get("title") or Path(filename).stem

    if ext not in ALLOWED_EXT:
        return jsonify({
            "error": f"Unsupported file type '{ext}'. Allowed: {', '.join(ALLOWED_EXT)}"
        }), 400

    file_bytes = f.read()
    if len(file_bytes) > MAX_FILE_MB * 1024 * 1024:
        return jsonify({"error": f"File too large. Max: {MAX_FILE_MB}MB"}), 413

    try:
        # 1. Extract text
        raw_text = _extract_text(file_bytes, filename)
        if not raw_text.strip():
            return jsonify({"error": "No text could be extracted from the file"}), 400

        # 2. Chunk
        chunks = _chunk_text(raw_text)
        if not chunks:
            return jsonify({"error": "Document produced no usable text chunks"}), 400

        # 3. Embed
        embedder   = _get_embedder()
        embeddings = embedder.encode(chunks).tolist()

        # 4. Store in ChromaDB
        doc_id    = str(uuid.uuid4())
        collection = _get_collection()

        chunk_ids = [f"{doc_id}_chunk_{i}" for i in range(len(chunks))]
        metadatas = [
            {
                "doc_id":    doc_id,
                "title":     title,
                "filename":  filename,
                "chunk_idx": i,
                "total_chunks": len(chunks),
                "uploaded_at": _now(),
            }
            for i in range(len(chunks))
        ]

        collection.add(
            ids=chunk_ids,
            documents=chunks,
            embeddings=embeddings,
            metadatas=metadatas,
        )

        logger.info(f"Ingested document '{title}' ({len(chunks)} chunks, doc_id={doc_id})")

        return jsonify({
            "doc_id":      doc_id,
            "title":       title,
            "filename":    filename,
            "chunks":      len(chunks),
            "chars":       len(raw_text),
            "message":     f"'{title}' ingested successfully ({len(chunks)} chunks)",
        })

    except (ImportError, ValueError) as e:
        return jsonify({"error": str(e)}), 500
    except Exception as e:
        logger.error(f"RAG upload error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": f"Ingestion failed: {str(e)}"}), 500


@rag_bp.route("/query", methods=["POST"])
def query_knowledge_base():
    """
    Query the RAG knowledge base and generate an AI answer with citations.
    Body (JSON):
      {
        "question":    str,          (required)
        "top_k":       int,          (optional, default 5)
        "model":       str,          (optional Ollama model override)
        "doc_filter":  str,          (optional doc_id to restrict search)
      }
    """
    data       = request.get_json(silent=True) or {}
    question   = (data.get("question") or "").strip()
    top_k      = min(int(data.get("top_k", TOP_K)), 10)
    model      = data.get("model")
    doc_filter = data.get("doc_filter")

    if not question:
        return jsonify({"error": "question is required"}), 400

    try:
        collection = _get_collection()
        if collection.count() == 0:
            return jsonify({
                "answer":      "No documents have been uploaded yet. Please upload a document first.",
                "citations":   [],
                "chunks_used": 0,
            })

        # Embed the question
        embedder   = _get_embedder()
        q_embedding = embedder.encode([question]).tolist()[0]

        # Build where clause for doc filter
        where = {"doc_id": doc_filter} if doc_filter else None

        # Retrieve relevant chunks
        results = collection.query(
            query_embeddings=[q_embedding],
            n_results=min(top_k, collection.count()),
            where=where,
            include=["documents", "metadatas", "distances"],
        )

        docs      = results.get("documents", [[]])[0]
        metas     = results.get("metadatas", [[]])[0]
        distances = results.get("distances", [[]])[0]

        if not docs:
            return jsonify({
                "answer":    "No relevant context found for your question.",
                "citations": [],
                "chunks_used": 0,
            })

        # Build context for LLM
        context_parts = []
        citations     = []
        for i, (doc, meta, dist) in enumerate(zip(docs, metas, distances)):
            relevance = round(1 - dist, 3)   # cosine similarity
            context_parts.append(
                f"[Source {i+1}: {meta.get('title','Unknown')}, Chunk {meta.get('chunk_idx',0)+1}]\n{doc}"
            )
            citations.append({
                "source_num": i + 1,
                "title":      meta.get("title", "Unknown"),
                "filename":   meta.get("filename", ""),
                "chunk_idx":  meta.get("chunk_idx", 0),
                "relevance":  relevance,
                "snippet":    doc[:200] + ("…" if len(doc) > 200 else ""),
            })

        context = "\n\n---\n\n".join(context_parts)

        # Generate answer with citations
        from .ollama_client import generate, OllamaUnavailableError

        prompt = (
            f"You are an expert business analyst answering questions based ONLY on the "
            f"provided document context. Always cite your sources using [Source N] notation.\n\n"
            f"Context from uploaded documents:\n{context}\n\n"
            f"Question: {question}\n\n"
            f"Answer (cite sources using [Source N]):"
        )

        try:
            answer = generate(prompt=prompt, model=model, temperature=0.3, max_tokens=1000)
        except OllamaUnavailableError:
            # Fallback: return raw chunks without AI synthesis
            answer = (
                "⚠️ Ollama is not available. Here are the most relevant excerpts from your documents:\n\n"
                + "\n\n".join([f"**{c['title']} (Chunk {c['chunk_idx']+1}):**\n{c['snippet']}" for c in citations])
            )

        return jsonify({
            "answer":      answer,
            "citations":   citations,
            "chunks_used": len(docs),
            "question":    question,
        })

    except Exception as e:
        logger.error(f"RAG query error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": f"Query failed: {str(e)}"}), 500


@rag_bp.route("/documents", methods=["GET"])
def list_documents():
    """List all ingested documents with metadata."""
    try:
        collection = _get_collection()
        if collection.count() == 0:
            return jsonify({"documents": [], "total_chunks": 0})

        # Get all chunk metadata
        all_items = collection.get(include=["metadatas"])
        metas     = all_items.get("metadatas", [])

        # Deduplicate by doc_id
        seen     = {}
        for meta in metas:
            doc_id = meta.get("doc_id", "unknown")
            if doc_id not in seen:
                seen[doc_id] = {
                    "doc_id":      doc_id,
                    "title":       meta.get("title", "Unknown"),
                    "filename":    meta.get("filename", ""),
                    "uploaded_at": meta.get("uploaded_at", ""),
                    "chunks":      meta.get("total_chunks", 0),
                }

        return jsonify({
            "documents":   list(seen.values()),
            "total_docs":  len(seen),
            "total_chunks": collection.count(),
        })

    except Exception as e:
        logger.error(f"List documents error: {e}")
        return jsonify({"documents": [], "error": str(e)}), 500


@rag_bp.route("/documents/<doc_id>", methods=["DELETE"])
def delete_document(doc_id: str):
    """Remove all chunks for a specific document from ChromaDB."""
    try:
        collection = _get_collection()
        # Find chunk IDs belonging to this doc
        all_items  = collection.get(where={"doc_id": doc_id}, include=["metadatas"])
        chunk_ids  = all_items.get("ids", [])
        if not chunk_ids:
            return jsonify({"error": f"Document '{doc_id}' not found"}), 404
        collection.delete(ids=chunk_ids)
        logger.info(f"Deleted document {doc_id} ({len(chunk_ids)} chunks)")
        return jsonify({"deleted": doc_id, "chunks_removed": len(chunk_ids)})
    except Exception as e:
        logger.error(f"Delete document error: {e}")
        return jsonify({"error": str(e)}), 500


@rag_bp.route("/documents", methods=["DELETE"])
def clear_all_documents():
    """Remove ALL documents from the knowledge base."""
    global _chroma_collection
    try:
        import chromadb
        client = chromadb.PersistentClient(path=CHROMA_DIR)
        client.delete_collection(COLLECTION_NAME)
        _chroma_collection = None  # Reset singleton so it recreates
        logger.info("Cleared all RAG documents")
        return jsonify({"cleared": True, "message": "All documents removed from knowledge base"})
    except Exception as e:
        logger.error(f"Clear documents error: {e}")
        return jsonify({"error": str(e)}), 500


@rag_bp.route("/status", methods=["GET"])
def rag_status():
    """Return RAG system status."""
    try:
        collection = _get_collection()
        count      = collection.count()
        return jsonify({
            "ready":        True,
            "total_chunks": count,
            "chroma_dir":   CHROMA_DIR,
            "collection":   COLLECTION_NAME,
        })
    except Exception as e:
        return jsonify({"ready": False, "error": str(e)}), 503


# ── Helpers ───────────────────────────────────────────────────────────

def _now() -> str:
    return datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ")
